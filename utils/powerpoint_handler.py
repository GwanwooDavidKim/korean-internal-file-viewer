# -*- coding: utf-8 -*-
"""
PowerPoint íŒŒì¼ ì²˜ë¦¬ ëª¨ë“ˆ (PowerPoint File Handler)

python-pptxë¥¼ ì‚¬ìš©í•˜ì—¬ PowerPoint íŒŒì¼ì—ì„œ í…ìŠ¤íŠ¸ì™€ ìŠ¬ë¼ì´ë“œ ì •ë³´ë¥¼ ì¶”ì¶œí•©ë‹ˆë‹¤.
PDF ë³€í™˜ì„ í†µí•œ ì•ˆì „í•œ ë¯¸ë¦¬ë³´ê¸° ê¸°ëŠ¥ì„ ì œê³µí•©ë‹ˆë‹¤.

í•µì‹¬ ê°œì„ ì‚¬í•­:
- ì‚¬ìš©ìì˜ PowerPoint ì‘ì—…ì— ì˜í–¥ ì—†ìŒ (PDF ë³€í™˜ ë°©ì‹)
- ì›ë³¸ íŒŒì¼ ë½ ì—†ìŒ
- "ì›ë³¸ ì—´ê¸°" ê¸°ëŠ¥ ì™„ë²½ ì‘ë™
- ë¹ ë¥´ê³  ì•ˆì •ì ì¸ ë Œë”ë§
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from .ppt_to_pdf_converter import get_converter
from .com_powerpoint_converter import get_com_converter
from .pdf_handler import PdfHandler
import logging
import time

# PILì„ ì•ˆì „í•˜ê²Œ import
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: Pillow not installed. PowerPoint image preview will not be available.")

logger = logging.getLogger(__name__)


class PowerPointHandler:
    """
    PowerPoint íŒŒì¼ ì²˜ë¦¬ë¥¼ ìœ„í•œ í´ë˜ìŠ¤ì…ë‹ˆë‹¤.
    
    ì£¼ìš” ê¸°ëŠ¥:
    - PowerPoint ìŠ¬ë¼ì´ë“œ í…ìŠ¤íŠ¸ ì¶”ì¶œ
    - ì•ˆì „í•œ PDF ë³€í™˜ì„ í†µí•œ ë¯¸ë¦¬ë³´ê¸°
    - ìŠ¬ë¼ì´ë“œ êµ¬ì¡° ë¶„ì„
    - í”„ë ˆì  í…Œì´ì…˜ ë©”íƒ€ë°ì´í„° ì¡°íšŒ
    - ì‚¬ìš©ì ì‘ì—… ë°©í•´ ì—†ìŒ (ì™„ì „ ê²©ë¦¬)
    """
    
    def __init__(self):
        """PowerPointHandler ì¸ìŠ¤í„´ìŠ¤ë¥¼ ì´ˆê¸°í™”í•©ë‹ˆë‹¤."""
        # PDF ë³€í™˜ê¸°ì™€ PDF í•¸ë“¤ëŸ¬ ì´ˆê¸°í™” (ë¨¼ì € ìƒì„±)
        # Windows+Office í™˜ê²½ì—ì„œëŠ” COM ë°©ì‹ ìš°ì„  ì‚¬ìš©
        self.com_converter = get_com_converter()
        self.pdf_converter = get_converter()  # í´ë°±ìš©
        self.pdf_handler = PdfHandler()
        
        # ì‚¬ìš©í•  ë³€í™˜ê¸° ê²°ì •
        if self.com_converter.is_available():
            self.active_converter = self.com_converter
            self.converter_type = "COM"
            self.supported_extensions = ['.ppt', '.pptx']  # COMì€ ëª¨ë“  PowerPoint í˜•ì‹ ì§€ì›
            print("   ğŸš€ Microsoft Office COM ë°©ì‹ ì‚¬ìš© (ê³ ì„±ëŠ¥)")
            print("   ğŸ“„ ì§€ì› í˜•ì‹: .ppt, .pptx")
        else:
            self.active_converter = self.pdf_converter
            self.converter_type = "LibreOffice" 
            self.supported_extensions = ['.pptx']  # LibreOfficeëŠ” .pptxë§Œ ì•ˆì •ì 
            print("   ğŸ“‹ LibreOffice ë°©ì‹ ì‚¬ìš© (í˜¸í™˜ì„±)")
            print("   ğŸ“„ ì§€ì› í˜•ì‹: .pptx")
        
        # í˜„ì¬ ì—°ê²°ëœ íŒŒì¼ ê²½ë¡œ (í˜¸í™˜ì„±ì„ ìœ„í•´)
        self.current_file_path = None
        
        print("ğŸ”„ PowerPointHandler ì´ˆê¸°í™” - ì•ˆì „í•œ PDF ë³€í™˜ ë°©ì‹ ì‚¬ìš©")
        print("   âœ… ì‚¬ìš©ì PowerPoint ì‘ì—…ì— ì˜í–¥ ì—†ìŒ")
        print("   âœ… ì›ë³¸ íŒŒì¼ ë½ ì—†ìŒ") 
        print("   âœ… 'ì›ë³¸ ì—´ê¸°' ê¸°ëŠ¥ ì™„ë²½ ì‘ë™")
        print(f"   âš¡ í™œì„± ë³€í™˜ê¸°: {self.converter_type}")
    
    def open_persistent_connection(self, file_path: str) -> bool:
        """
        í˜¸í™˜ì„±ì„ ìœ„í•œ ë©”ì†Œë“œ - PDF ë³€í™˜ ë°©ì‹ì—ì„œëŠ” ì§€ì† ì—°ê²°ì´ ë¶ˆí•„ìš”
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            
        Returns:
            bool: í•­ìƒ True (PDF ë³€í™˜ ë°©ì‹ì€ í•­ìƒ ì‚¬ìš© ê°€ëŠ¥)
        """
        self.current_file_path = file_path  # í˜„ì¬ íŒŒì¼ ê²½ë¡œ ì €ì¥ (render_slide_fastìš©)
        logger.info(f"ğŸ”„ PPT â†’ PDF ë°©ì‹ìœ¼ë¡œ ì—°ê²°: {os.path.basename(file_path)}")
        logger.info("   âœ… ì§€ì† ì—°ê²° ë¶ˆí•„ìš” - ì¦‰ì‹œ ë Œë”ë§ ê°€ëŠ¥")
        return True
    
    def close_persistent_connection(self):
        """
        í˜¸í™˜ì„±ì„ ìœ„í•œ ë©”ì†Œë“œ - PDF ë³€í™˜ ë°©ì‹ì—ì„œëŠ” ì •ë¦¬í•  ì—°ê²°ì´ ì—†ìŒ
        """
        self.current_file_path = None  # í˜„ì¬ íŒŒì¼ ê²½ë¡œ ì´ˆê¸°í™”
        logger.info("ğŸ”„ PPT â†’ PDF ë°©ì‹ ì •ë¦¬ ì™„ë£Œ")
        logger.info("   âœ… ì‚¬ìš©ì PowerPointì— ì˜í–¥ ì—†ì´ ì•ˆì „í•˜ê²Œ ì¢…ë£Œ")
    
    def is_connected(self) -> bool:
        """
        í˜¸í™˜ì„±ì„ ìœ„í•œ ë©”ì†Œë“œ - ì—°ê²° ìƒíƒœ í™•ì¸
        
        Returns:
            bool: íŒŒì¼ì´ ì—°ê²°ë˜ì–´ ìˆëŠ”ì§€ ì—¬ë¶€
        """
        return self.current_file_path is not None
    
    def render_slide_fast(self, slide_number: int, width: int = 800, height: int = 600) -> Optional['Image.Image']:
        """
        í˜¸í™˜ì„±ì„ ìœ„í•œ ë©”ì†Œë“œ - PDF ë³€í™˜ ë°©ì‹ì—ì„œëŠ” ë¹ ë¥¸/ì¼ë°˜ ë Œë”ë§ êµ¬ë¶„ì´ ì—†ìŒ
        
        Args:
            slide_number (int): ìŠ¬ë¼ì´ë“œ ë²ˆí˜¸ (0ë¶€í„° ì‹œì‘)
            width (int): ì´ë¯¸ì§€ ë„ˆë¹„
            height (int): ì´ë¯¸ì§€ ë†’ì´
            
        Returns:
            Optional[Image.Image]: ë Œë”ë§ëœ ì´ë¯¸ì§€
        """
        if not self.current_file_path:
            logger.error("âŒ render_slide_fast í˜¸ì¶œ ì „ì— open_persistent_connectionì´ í•„ìš”í•©ë‹ˆë‹¤")
            return None
            
        logger.info(f"ğŸš€ ë¹ ë¥¸ ë Œë”ë§ (PDF ë°©ì‹): ìŠ¬ë¼ì´ë“œ {slide_number + 1}")
        # PDF ë³€í™˜ ë°©ì‹ì€ í•­ìƒ ë¹ ë¥´ë¯€ë¡œ ê¸°ë³¸ ë Œë”ë§ ë©”ì†Œë“œì™€ ë™ì¼
        return self.render_slide_to_image(self.current_file_path, slide_number, width, height)
    
    def can_handle(self, file_path: str) -> bool:
        """
        íŒŒì¼ì´ ì´ í•¸ë“¤ëŸ¬ê°€ ì²˜ë¦¬í•  ìˆ˜ ìˆëŠ” í˜•ì‹ì¸ì§€ í™•ì¸í•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): íŒŒì¼ ê²½ë¡œ
            
        Returns:
            bool: ì²˜ë¦¬ ê°€ëŠ¥ ì—¬ë¶€
        """
        return any(file_path.lower().endswith(ext) for ext in self.supported_extensions)
    
    def get_slide_count(self, file_path: str) -> int:
        """
        PowerPointì˜ ì´ ìŠ¬ë¼ì´ë“œ ìˆ˜ë¥¼ ë°˜í™˜í•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            
        Returns:
            int: ìŠ¬ë¼ì´ë“œ ìˆ˜ (ì˜¤ë¥˜ ì‹œ 0)
        """
        try:
            prs = Presentation(file_path)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"ìŠ¬ë¼ì´ë“œ ìˆ˜ í™•ì¸ ì˜¤ë¥˜: {e}")
            return 0
    
    def render_slide_to_image(self, file_path: str, slide_number: int, width: int = 800, height: int = 600) -> Optional['Image.Image']:
        """
        PPT â†’ PDF â†’ ì´ë¯¸ì§€ ë°©ì‹ìœ¼ë¡œ ì•ˆì „í•˜ê²Œ ìŠ¬ë¼ì´ë“œë¥¼ ë Œë”ë§í•©ë‹ˆë‹¤.
        
        ì´ ë°©ì‹ì˜ ì¥ì :
        - ì‚¬ìš©ìì˜ PowerPoint ì‘ì—…ì— ì˜í–¥ ì—†ìŒ
        - ì›ë³¸ íŒŒì¼ì´ ì ê¸°ì§€ ì•ŠìŒ  
        - "ì›ë³¸ ì—´ê¸°" ê¸°ëŠ¥ ì™„ë²½ ì‘ë™
        - ë¹ ë¥´ê³  ì•ˆì •ì ì¸ ë Œë”ë§
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            slide_number (int): ìŠ¬ë¼ì´ë“œ ë²ˆí˜¸ (0ë¶€í„° ì‹œì‘)
            width (int): ì´ë¯¸ì§€ ë„ˆë¹„ (ë¯¸ì‚¬ìš© - PDF ê¸°ë³¸ í•´ìƒë„)
            height (int): ì´ë¯¸ì§€ ë†’ì´ (ë¯¸ì‚¬ìš© - PDF ê¸°ë³¸ í•´ìƒë„)
            
        Returns:
            Optional[Image.Image]: ìƒì„±ëœ ì´ë¯¸ì§€ (ì‹¤íŒ¨ ì‹œ None)
        """
        if not PIL_AVAILABLE:
            logger.warning("PILì´ ì—†ì–´ì„œ ì´ë¯¸ì§€ ë Œë”ë§ ë¶ˆê°€")
            return None
        
        try:
            logger.info(f"ğŸ”„ PPT â†’ PDF â†’ ì´ë¯¸ì§€ ë Œë”ë§ ì‹œì‘: {os.path.basename(file_path)}, ìŠ¬ë¼ì´ë“œ {slide_number + 1}")
            
            # 1ë‹¨ê³„: PPTë¥¼ PDFë¡œ ë³€í™˜ (ìºì‹œ í™œìš©) - í™œì„± ë³€í™˜ê¸° ì‚¬ìš©
            start_time = time.time()
            pdf_path = self.active_converter.convert_to_pdf(file_path)
            conversion_time = time.time() - start_time
            if not pdf_path:
                logger.error("âŒ PPT â†’ PDF ë³€í™˜ ì‹¤íŒ¨")
                return None
            
            logger.info(f"âœ… PDF ë³€í™˜ ì™„ë£Œ: {os.path.basename(pdf_path)}")
            
            # 2ë‹¨ê³„: PDFì—ì„œ í•´ë‹¹ í˜ì´ì§€ë¥¼ ì´ë¯¸ì§€ë¡œ ë Œë”ë§
            image = self.pdf_handler.render_page_to_image(
                pdf_path, 
                page_num=slide_number,
                zoom=1.5  # ê³ í’ˆì§ˆì„ ìœ„í•œ 150% í™•ëŒ€
            )
            
            if image:
                logger.info(f"âœ… ìŠ¬ë¼ì´ë“œ {slide_number + 1} ë Œë”ë§ ì™„ë£Œ! ({self.converter_type} ë³€í™˜: {conversion_time:.1f}ì´ˆ)")
                return image
            else:
                logger.error(f"âŒ PDF í˜ì´ì§€ {slide_number} ë Œë”ë§ ì‹¤íŒ¨")
                return None
                
        except Exception as e:
            logger.error(f"âŒ ìŠ¬ë¼ì´ë“œ ë Œë”ë§ ì˜¤ë¥˜: {e}")
            return None
    
    def extract_text_from_slide(self, file_path: str, slide_number: int) -> Dict[str, Any]:
        """
        ì§€ì •ëœ ìŠ¬ë¼ì´ë“œì—ì„œ í…ìŠ¤íŠ¸ë¥¼ ì¶”ì¶œí•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            slide_number (int): ìŠ¬ë¼ì´ë“œ ë²ˆí˜¸ (0ë¶€í„° ì‹œì‘)
            
        Returns:
            Dict[str, Any]: ìŠ¬ë¼ì´ë“œ ì •ë³´
        """
        try:
            prs = Presentation(file_path)
            
            if slide_number >= len(prs.slides) or slide_number < 0:
                return {'error': 'ì˜ëª»ëœ ìŠ¬ë¼ì´ë“œ ë²ˆí˜¸'}
            
            slide = prs.slides[slide_number]
            
            # ìŠ¬ë¼ì´ë“œ ì œëª© ì¶”ì¶œ
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            
            # í…ìŠ¤íŠ¸ ë‚´ìš© ì¶”ì¶œ
            text_content = []
            bullet_points = []
            
            for shape in slide.shapes:
                # í…ìŠ¤íŠ¸ê°€ ìˆëŠ” shapeë§Œ ì²˜ë¦¬
                if hasattr(shape, "text") and hasattr(shape, "text_frame") and shape.text.strip():
                    # ì œëª©ì´ ì•„ë‹Œ ê²½ìš°ì—ë§Œ ì¶”ê°€
                    if shape != slide.shapes.title:
                        text_content.append(shape.text)
                        
                        # í…ìŠ¤íŠ¸ í”„ë ˆì„ì´ ìˆëŠ” ê²½ìš° ë‹¨ë½ë³„ë¡œ ë¶„ì„
                        if shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    bullet_points.append({
                                        'text': paragraph.text,
                                        'level': paragraph.level,
                                    })
            
            # ì´ë¯¸ì§€ ë° ê¸°íƒ€ ê°ì²´ ì¹´ìš´íŠ¸
            image_count = 0
            chart_count = 0
            table_count = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_count += 1
            
            return {
                'slide_number': slide_number + 1,
                'title': title or '[ì œëª© ì—†ìŒ]',
                'text_content': text_content,
                'bullet_points': bullet_points,
                'full_text': "\n".join(text_content),
                'image_count': image_count,
                'chart_count': chart_count,
                'table_count': table_count,
                'total_shapes': len(slide.shapes),
            }
            
        except Exception as e:
            logger.error(f"ìŠ¬ë¼ì´ë“œ í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
            return {'error': f"ìŠ¬ë¼ì´ë“œ í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}"}
    
    def extract_text(self, file_path: str, max_slides: int = None) -> str:
        """
        ì „ì²´ í”„ë ˆì  í…Œì´ì…˜ì—ì„œ í…ìŠ¤íŠ¸ë¥¼ ì¶”ì¶œí•©ë‹ˆë‹¤.
        (ê²€ìƒ‰ ì¸ë±ì‹±ìš©)
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            max_slides (int, optional): ìµœëŒ€ ìŠ¬ë¼ì´ë“œ ìˆ˜ ì œí•œ (Noneì´ë©´ ëª¨ë“  ìŠ¬ë¼ì´ë“œ)
            
        Returns:
            str: ì¶”ì¶œëœ ì „ì²´ í…ìŠ¤íŠ¸
        """
        # .ppt íŒŒì¼ì€ python-pptxë¡œ ì§ì ‘ ì½ì„ ìˆ˜ ì—†ìœ¼ë¯€ë¡œ PDFì—ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ
        if file_path.lower().endswith('.ppt'):
            try:
                logger.info(f"ğŸ”„ .ppt íŒŒì¼ í…ìŠ¤íŠ¸ ì¶”ì¶œ: PDF ë³€í™˜ ë°©ì‹ ì‚¬ìš©")
                if max_slides is not None:
                    logger.info(f"í…ìŠ¤íŠ¸ ì¶”ì¶œ ì œí•œ: ìµœëŒ€ {max_slides}ê°œ ìŠ¬ë¼ì´ë“œë§Œ ì²˜ë¦¬")
                pdf_path = self.active_converter.convert_to_pdf(file_path)
                if pdf_path:
                    # PDFì—ì„œ max_slidesë¥¼ max_pagesë¡œ ë³€í™˜í•˜ì—¬ ì „ë‹¬
                    return self.pdf_handler.extract_text(pdf_path, max_pages=max_slides)
                else:
                    return f".ppt íŒŒì¼ í…ìŠ¤íŠ¸ ì¶”ì¶œ ì‹¤íŒ¨: {os.path.basename(file_path)}"
            except Exception as e:
                logger.error(f".ppt í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
                return f".ppt íŒŒì¼ í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}"
        
        # .pptx íŒŒì¼ì€ python-pptxë¡œ ì§ì ‘ ì¶”ì¶œ
        try:
            prs = Presentation(file_path)
            all_text = []
            
            # ìŠ¬ë¼ì´ë“œ ìˆ˜ ì œí•œ ì ìš©
            slides_to_process = prs.slides
            if max_slides is not None:
                slides_to_process = list(prs.slides)[:max_slides]
                logger.info(f"í…ìŠ¤íŠ¸ ì¶”ì¶œ ì œí•œ: ìµœëŒ€ {max_slides}ê°œ ìŠ¬ë¼ì´ë“œë§Œ ì²˜ë¦¬")
            
            for i, slide in enumerate(slides_to_process):
                slide_text = []
                
                # ìŠ¬ë¼ì´ë“œ ì œëª©
                if slide.shapes.title:
                    slide_text.append(f"=== ìŠ¬ë¼ì´ë“œ {i + 1}: {slide.shapes.title.text} ===")
                else:
                    slide_text.append(f"=== ìŠ¬ë¼ì´ë“œ {i + 1} ===")
                
                # ìŠ¬ë¼ì´ë“œ ë‚´ìš©
                for shape in slide.shapes:
                    if hasattr(shape, "text") and hasattr(shape, "text_frame") and shape.text.strip():
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text)
                
                all_text.append("\n".join(slide_text))
            
            return "\n\n".join(all_text)
            
        except Exception as e:
            logger.error(f"PowerPoint í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
            return f"PowerPoint í…ìŠ¤íŠ¸ ì¶”ì¶œ ì˜¤ë¥˜: {e}"
    
    def extract_all_text(self, file_path: str, max_slides: int = None) -> str:
        """extract_textì˜ ë³„ì¹­ (í˜¸í™˜ì„±ì„ ìœ„í•´)"""
        return self.extract_text(file_path, max_slides)
    
    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """
        í”„ë ˆì  í…Œì´ì…˜ì˜ ìƒì„¸ ì •ë³´ë¥¼ ë°˜í™˜í•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            
        Returns:
            Dict[str, Any]: í”„ë ˆì  í…Œì´ì…˜ ì •ë³´
        """
        try:
            if not os.path.exists(file_path):
                return {'error': 'íŒŒì¼ì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤'}
            
            # .ppt íŒŒì¼ì€ python-pptxë¡œ ì½ì„ ìˆ˜ ì—†ìœ¼ë¯€ë¡œ PDF ê¸°ë°˜ ì •ë³´ ì¶”ì¶œ
            if file_path.lower().endswith('.ppt'):
                return self._get_ppt_info_via_pdf(file_path)
            
            # .pptx íŒŒì¼ì€ python-pptxë¡œ ì§ì ‘ ì²˜ë¦¬
            prs = Presentation(file_path)
            
            # ê¸°ë³¸ ì •ë³´
            file_size = os.path.getsize(file_path)
            slide_count = len(prs.slides)
            
            # í”„ë ˆì  í…Œì´ì…˜ ë©”íƒ€ë°ì´í„°
            core_props = prs.core_properties
            
            # ìŠ¬ë¼ì´ë“œ í¬ê¸° ì •ë³´
            slide_width = prs.slide_width or 9144000  # ê¸°ë³¸ê°’ (10ì¸ì¹˜)
            slide_height = prs.slide_height or 6858000  # ê¸°ë³¸ê°’ (7.5ì¸ì¹˜)
            
            # ê° ìŠ¬ë¼ì´ë“œì˜ ìš”ì•½ ì •ë³´
            slides_summary = []
            total_images = 0
            total_charts = 0
            total_tables = 0
            
            for i, slide in enumerate(prs.slides):
                # ìŠ¬ë¼ì´ë“œ ì œëª©
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text
                
                # ê°ì²´ ì¹´ìš´íŠ¸
                images = charts = tables = 0
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        charts += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        tables += 1
                
                total_images += images
                total_charts += charts
                total_tables += tables
                
                # í…ìŠ¤íŠ¸ ë¸”ë¡ ìˆ˜
                text_shapes = sum(1 for shape in slide.shapes 
                                if hasattr(shape, "text") and hasattr(shape, "text_frame") and shape.text.strip())
                
                slides_summary.append({
                    'slide_number': i + 1,
                    'title': title or f'[ìŠ¬ë¼ì´ë“œ {i + 1}]',
                    'text_shapes': text_shapes,
                    'images': images,
                    'charts': charts,
                    'tables': tables,
                    'total_shapes': len(slide.shapes),
                })
            
            # PDF ë³€í™˜ ê°€ëŠ¥ ì—¬ë¶€ í™•ì¸ (í™œì„± ë³€í™˜ê¸°ì—ì„œ ì§ì ‘)
            conversion_info = self.active_converter.get_cache_info()
            conversion_available = conversion_info.get('converter_available', 
                                                      conversion_info.get('libreoffice_available', False))
            
            info = {
                'filename': os.path.basename(file_path),
                'file_size': file_size,
                'file_size_mb': round(file_size / (1024 * 1024), 2),
                'slide_count': slide_count,
                'slide_width_inches': round(slide_width / 914400, 2),
                'slide_height_inches': round(slide_height / 914400, 2),
                'total_images': total_images,
                'total_charts': total_charts,
                'total_tables': total_tables,
                'slides_summary': slides_summary,
                'conversion_available': conversion_available,
                'converter_type': self.converter_type,
                'cache_info': conversion_info,
                'metadata': {
                    'title': getattr(core_props, 'title', None),
                    'subject': getattr(core_props, 'subject', None),
                    'author': getattr(core_props, 'author', None),
                    'created': getattr(core_props, 'created', None),
                    'last_modified_by': getattr(core_props, 'last_modified_by', None),
                    'modified': getattr(core_props, 'modified', None),
                }
            }
            
            return info
            
        except Exception as e:
            logger.error(f"í”„ë ˆì  í…Œì´ì…˜ ì •ë³´ ì¡°íšŒ ì˜¤ë¥˜: {e}")
            return {'error': f'í”„ë ˆì  í…Œì´ì…˜ ì •ë³´ ì¡°íšŒ ì˜¤ë¥˜: {e}'}
    
    def _get_ppt_info_via_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        .ppt íŒŒì¼ì˜ ì •ë³´ë¥¼ PDF ë³€í™˜ì„ í†µí•´ ì¶”ì¶œí•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): .ppt íŒŒì¼ ê²½ë¡œ
            
        Returns:
            Dict[str, Any]: í”„ë ˆì  í…Œì´ì…˜ ì •ë³´
        """
        try:
            logger.info(f"ğŸ”„ .ppt íŒŒì¼ ì •ë³´ ì¶”ì¶œ: PDF ë³€í™˜ ë°©ì‹ ì‚¬ìš©")
            
            # ê¸°ë³¸ íŒŒì¼ ì •ë³´
            file_size = os.path.getsize(file_path)
            
            # PDFë¡œ ë³€í™˜í•˜ì—¬ ìŠ¬ë¼ì´ë“œ ìˆ˜ í™•ì¸
            pdf_path = self.active_converter.convert_to_pdf(file_path)
            slide_count = 0
            if pdf_path:
                # PDF í•¸ë“¤ëŸ¬ë¡œ í˜ì´ì§€ ìˆ˜ í™•ì¸ (= ìŠ¬ë¼ì´ë“œ ìˆ˜)
                doc = self.pdf_handler._open_document(pdf_path)
                if doc:
                    slide_count = len(doc)
                    doc.close()
            
            # ë³€í™˜ ì •ë³´
            conversion_info = self.active_converter.get_cache_info()
            conversion_available = conversion_info.get('converter_available', 
                                                      conversion_info.get('libreoffice_available', False))
            
            return {
                'filename': os.path.basename(file_path),
                'file_size': file_size,
                'file_size_mb': round(file_size / (1024 * 1024), 2),
                'slide_count': slide_count,
                'file_type': '.ppt (Legacy PowerPoint)',
                'conversion_available': conversion_available,
                'converter_type': self.converter_type,
                'cache_info': conversion_info,
                'note': '.ppt íŒŒì¼ì€ PDF ë³€í™˜ì„ í†µí•´ì„œë§Œ ë¯¸ë¦¬ë³´ê¸° ê°€ëŠ¥í•©ë‹ˆë‹¤.',
                'slides_summary': [{'slide_number': i+1, 'title': f'ìŠ¬ë¼ì´ë“œ {i+1}', 'note': 'PDF ë³€í™˜ ë°©ì‹'} 
                                  for i in range(slide_count)],
                'metadata': {
                    'title': None,
                    'subject': None, 
                    'author': None,
                    'created': None,
                    'last_modified_by': None,
                    'modified': None,
                }
            }
            
        except Exception as e:
            logger.error(f".ppt íŒŒì¼ ì •ë³´ ì¶”ì¶œ ì˜¤ë¥˜: {e}")
            return {'error': f'.ppt íŒŒì¼ ì •ë³´ ì¶”ì¶œ ì˜¤ë¥˜: {e}'}
    
    def search_in_presentation(self, file_path: str, search_term: str, max_results: int = 10) -> List[Dict[str, Any]]:
        """
        í”„ë ˆì  í…Œì´ì…˜ì—ì„œ í…ìŠ¤íŠ¸ë¥¼ ê²€ìƒ‰í•©ë‹ˆë‹¤.
        
        Args:
            file_path (str): PowerPoint íŒŒì¼ ê²½ë¡œ
            search_term (str): ê²€ìƒ‰ì–´
            max_results (int): ìµœëŒ€ ê²°ê³¼ ìˆ˜
            
        Returns:
            List[Dict[str, Any]]: ê²€ìƒ‰ ê²°ê³¼
        """
        try:
            prs = Presentation(file_path)
            results = []
            search_term_lower = search_term.lower()
            
            for i, slide in enumerate(prs.slides):
                # ìŠ¬ë¼ì´ë“œ ì œëª© ê²€ìƒ‰
                if slide.shapes.title and slide.shapes.title.text:
                    title_text = slide.shapes.title.text
                    if search_term_lower in title_text.lower():
                        results.append({
                            'slide_number': i + 1,
                            'location': 'ì œëª©',
                            'type': 'title',
                            'text': title_text,
                            'context': title_text,
                        })
                        
                        if len(results) >= max_results:
                            return results
                
                # ìŠ¬ë¼ì´ë“œ ë‚´ìš© ê²€ìƒ‰
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and hasattr(shape, "text_frame") and shape.text.strip():
                        if shape != slide.shapes.title:
                            text = shape.text
                            if search_term_lower in text.lower():
                                # ì»¨í…ìŠ¤íŠ¸ ìƒì„± (ê²€ìƒ‰ì–´ ì£¼ë³€ í…ìŠ¤íŠ¸)
                                context_start = max(0, text.lower().find(search_term_lower) - 50)
                                context_end = min(len(text), text.lower().find(search_term_lower) + len(search_term) + 50)
                                context = text[context_start:context_end]
                                
                                if context_start > 0:
                                    context = "..." + context
                                if context_end < len(text):
                                    context = context + "..."
                                
                                results.append({
                                    'slide_number': i + 1,
                                    'location': f'í…ìŠ¤íŠ¸ ë¸”ë¡ {shape_idx + 1}',
                                    'type': 'content',
                                    'text': text,
                                    'context': context,
                                })
                                
                                if len(results) >= max_results:
                                    return results
            
            return results
            
        except Exception as e:
            logger.error(f"í”„ë ˆì  í…Œì´ì…˜ ê²€ìƒ‰ ì˜¤ë¥˜: {e}")
            return [{'error': f"í”„ë ˆì  í…Œì´ì…˜ ê²€ìƒ‰ ì˜¤ë¥˜: {e}"}]